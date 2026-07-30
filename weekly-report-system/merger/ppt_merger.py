"""
PPT 合并器（海外 BD）
将多人 PPT 统一字体后合并为单一 PPT，每人幻灯片前插入分隔页
"""
import os
import re
import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


class PptMerger:
    """PPT 合并器 - 海外 BD 团队周报合并"""

    CN_FONT = "等线"      # 等线
    EN_FONT = "Arial"
    FONT_SIZE_BODY = Pt(12)

    def __init__(self, config_path="config.yaml"):
        """
        初始化 PPT 合并器，从 config.yaml 读取配置

        参数:
            config_path: 配置文件路径（相对于项目根目录或当前工作目录）
        """
        # 解析配置文件路径（与 ExcelMerger 一致的逻辑）
        project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        config_full_path = os.path.join(project_root, config_path)
        if not os.path.exists(config_full_path):
            config_full_path = config_path  # fallback: 当前工作目录
        with open(config_full_path, "r", encoding="utf-8") as f:
            self.config = yaml.safe_load(f)

    def merge(self, upload_dir, output_dir, week_start, week_end):
        """
        扫描 upload_dir 中的 .pptx 文件，统一字体后合并。
        返回标准化摘要 dict。

        参数:
            upload_dir:  上传文件所在目录
            output_dir:  合并后文件输出目录
            week_start:  周起始日期 (YYYY-MM-DD)
            week_end:    周结束日期 (YYYY-MM-DD)

        返回:
            dict: {
                module_name, total_people, submitted_count,
                submission_rate, risk_items, key_projects,
                output_file, deadline_passed
            }
            若没有找到 pptx 文件则返回 None
        """
        # 扫描 .pptx 文件（排除临时文件 ~$）
        all_files = os.listdir(upload_dir) if os.path.isdir(upload_dir) else []
        ppt_files = sorted([
            f for f in all_files
            if f.endswith((".pptx", ".ppt")) and not f.startswith("~$")
        ])

        if not ppt_files:
            print(f"[提示] 在 {upload_dir} 中未找到 PPT 文件")
            return None

        print(f"找到 {len(ppt_files)} 个 PPT 文件，开始合并...")

        # 创建空白目标演示文稿
        merged = Presentation()

        for fname in ppt_files:
            person_name = self._extract_person_name(fname)

            # 插入分隔页
            self._add_divider_slide(merged, person_name)

            # 复制该人的所有幻灯片
            file_path = os.path.join(upload_dir, fname)
            try:
                src_prs = Presentation(file_path)
                for src_slide in src_prs.slides:
                    self._copy_slide(merged, src_slide)
            except Exception as e:
                print(f"[警告] 复制 {fname} 幻灯片失败: {e}")
                continue

        # 统一字体
        self._unify_fonts(merged)

        # 保存输出文件
        os.makedirs(output_dir, exist_ok=True)
        # 文件名中的日期使用 week_end（与 ExcelMerger 命名风格一致）
        output_path = os.path.join(output_dir, f"海外BD_周报汇总_{week_end}.pptx")
        merged.save(output_path)
        print(f"合并完成: 共 {len(ppt_files)} 人，输出 {output_path}")

        return self._build_summary(ppt_files, output_path, week_start, week_end)

    def _extract_person_name(self, filename):
        """
        从文件名提取姓名

        支持模式:
        - "张三_周报.pptx" → 张三
        - "weekly-report-李四.pptx" → 李四
        - "王五.pptx" → 王五
        - "David-Chen.pptx" → David-Chen（无中文时取文件名主体）
        """
        # 去掉扩展名
        name = os.path.splitext(filename)[0]

        # 清理常见前缀/后缀
        for prefix in ["周报_", "周报-", "weekly_", "weekly-", "Weekly_", "Weekly-"]:
            if name.startswith(prefix):
                name = name[len(prefix):]
        for suffix in ["_周报", "-周报", "_weekly", "-weekly", "_Weekly", "-Weekly"]:
            if name.endswith(suffix):
                name = name[:-len(suffix)]

        # 匹配中文姓名（2-4个汉字）
        m = re.search(r'([一-鿿]{2,4})', name)
        if m:
            return m.group(1)

        # 无中文时返回清理后的文件名主体
        return name.strip()

    def _add_divider_slide(self, prs, person_name):
        """
        插入分隔页（空白布局 + 居中人名标题）

        参数:
            prs:         目标 Presentation 对象
            person_name: 人员姓名
        """
        # 使用空白布局（索引 6）
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)

        # 居中放置文本框
        left = Inches(1)
        top = Inches(2.5)
        width = Inches(8)
        height = Inches(1.5)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = person_name
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.name = self.CN_FONT
        p.alignment = PP_ALIGN.CENTER

    def _copy_slide(self, dest_prs, src_slide):
        """
        复制幻灯片内容到目标演示文稿

        当前实现：复制文本框内容（周报 PPT 以文本为主）
        注意：图片、表格、图表等非文本元素不会复制

        参数:
            dest_prs:   目标 Presentation
            src_slide:  源 Slide 对象
        """
        # 使用空白布局
        layout = dest_prs.slide_layouts[6]
        dest_slide = dest_prs.slides.add_slide(layout)

        for shape in src_slide.shapes:
            if shape.has_text_frame:
                # 复制文本框
                txBox = dest_slide.shapes.add_textbox(
                    shape.left, shape.top, shape.width, shape.height
                )
                tf = txBox.text_frame
                tf.word_wrap = shape.text_frame.word_wrap

                for i, src_para in enumerate(shape.text_frame.paragraphs):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()

                    # 复制段落文本和格式
                    p.text = src_para.text
                    p.alignment = src_para.alignment if src_para.alignment is not None else PP_ALIGN.LEFT

                    # 复制 run 级别格式
                    for j, src_run in enumerate(src_para.runs):
                        if j < len(p.runs):
                            dest_run = p.runs[j]
                            if src_run.font.size:
                                dest_run.font.size = src_run.font.size
                            if src_run.font.bold:
                                dest_run.font.bold = src_run.font.bold
                            if src_run.font.italic:
                                dest_run.font.italic = src_run.font.italic
                            if src_run.font.name:
                                dest_run.font.name = src_run.font.name

    def _unify_fonts(self, prs):
        """
        统一所有幻灯片字体
        规则：含中文的 run → 等线，纯英文/数字 run → Arial
        保留原有字号，不强制覆盖为统一大小
        """
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            # 检测是否包含中文字符（CJK 统一表意文字）
                            if re.search(r'[一-鿿]', run.text):
                                run.font.name = self.CN_FONT
                            else:
                                run.font.name = self.EN_FONT

    def _build_summary(self, files, output_path, week_start, week_end):
        """
        构建标准化摘要 dict

        参数:
            files:       pptx 文件名列表
            output_path: 合并后的输出文件路径
            week_start:  周起始日期
            week_end:    周结束日期

        返回:
            dict: 标准化摘要
        """
        return {
            "module_name": "海外BD",
            "total_people": len(files),
            "submitted_count": len(files),
            "submission_rate": f"{len(files)}/{len(files)}",
            "risk_items": [],
            "key_projects": [],
            "output_file": output_path,
            "deadline_passed": False,
        }


# ===== 调试入口 =====
if __name__ == "__main__":
    import sys

    # 支持命令行测试: python ppt_merger.py <upload_dir> [output_dir]
    if len(sys.argv) >= 2:
        upload_dir = sys.argv[1]
        output_dir = sys.argv[2] if len(sys.argv) >= 3 else "output"
        week_start = sys.argv[3] if len(sys.argv) >= 4 else "2026-01-01"
        week_end = sys.argv[4] if len(sys.argv) >= 5 else "2026-01-07"
    else:
        # 默认测试：扫描 data/raw 目录
        project_root = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        upload_dir = os.path.join(project_root, "data", "raw")
        output_dir = os.path.join(project_root, "output")
        week_start = "2026-01-01"
        week_end = "2026-01-07"

    merger = PptMerger()
    result = merger.merge(upload_dir, output_dir, week_start, week_end)
    if result:
        print(f"\n摘要: {result}")
    else:
        print("\n未找到可合并的文件，创建测试 PPT 以供验证...")
        # 创建一个测试 PPT 文件
        test_dir = os.path.join(project_root, "data", "raw") if 'project_root' in dir() else upload_dir
        os.makedirs(test_dir, exist_ok=True)
        test_path = os.path.join(test_dir, "测试1_张三_周报.pptx")
        prs = Presentation()
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "本周工作总结"
        slide.placeholders[1].text = "项目A已完成上线，下周计划优化性能"
        prs.save(test_path)
        print(f"已创建测试文件: {test_path}")

        merger2 = PptMerger()
        result2 = merger2.merge(test_dir, output_dir, week_start, week_end)
        if result2:
            print(f"\n摘要: {result2}")
